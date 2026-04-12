"""
Content Indexing Service - Multi-format ingestion into ChromaDB.
Supports: PDF, Video (MP4/MOV), YouTube URL, PPTX, TXT, MD.
Pipeline: Extract text -> Chunk -> OpenAI Embed -> ChromaDB.
"""
import os
import uuid
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any
from datetime import datetime

logger = logging.getLogger(__name__)


class ContentIndexingService:
    CHUNK_SIZE = 800
    CHUNK_OVERLAP = 100
    # Maximum raw upload size in bytes (tune based on infra limits)
    MAX_FILE_SIZE = 500 * 1024 * 1024  # 500 MB — effectively unlimited  # 50MB
    EMBEDDING_MODEL = "text-embedding-3-small"

    def __init__(self, openai_client=None, chroma_client=None,
                 storage_path: str = "storage/documents",
                 whisper_model_name: str = None,
                 whisper_model=None):
        self.openai_client = openai_client
        self.chroma_client = chroma_client
        self.storage_path = Path(storage_path)
        self.storage_path.mkdir(parents=True, exist_ok=True)
        self._whisper_model = whisper_model  # injected from app.state (shared singleton)
        self._whisper_model_name = whisper_model_name or os.getenv("WHISPER_MODEL", "base")
        self.text_splitter = None
        try:
            from langchain.text_splitter import RecursiveCharacterTextSplitter
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.CHUNK_SIZE,
                chunk_overlap=self.CHUNK_OVERLAP,
                length_function=len,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
        except ImportError:
            logger.warning("langchain not available - using simple chunking")
        logger.info(f"ContentIndexingService initialized | whisper={self._whisper_model_name}")

    def _get_whisper_model(self):
        if self._whisper_model is not None:
            logger.info("Using pre-loaded Whisper model (shared singleton)")
            return self._whisper_model

        # Try to reuse the startup singleton from teaching_assistant_service
        try:
            from app.services.teaching_assistant_service import load_whisper_model
            logger.info("Reusing Whisper model from teaching_assistant_service singleton...")
            self._whisper_model = load_whisper_model()
            return self._whisper_model
        except Exception as e:
            logger.warning(f"Could not reuse TA whisper model: {e}")

        # Last resort: load our own copy
        try:
            import whisper
            import torch
            device = "cuda" if torch.cuda.is_available() else "cpu"
            logger.info(f"Loading Whisper {self._whisper_model_name} on {device}...")
            self._whisper_model = whisper.load_model(self._whisper_model_name, device=device)
        except ImportError:
            raise RuntimeError("Install whisper: pip install openai-whisper")
        return self._whisper_model

    def _get_chroma_collection(self, tenant_id: str, course_id: str):
        if not self.chroma_client:
            raise RuntimeError("ChromaDB client not configured")
        collection_name = f"tenant_{str(tenant_id)}_course_{str(course_id)}".replace("-", "_")
        return self.chroma_client.get_or_create_collection(
            name=collection_name,
            metadata={"tenant_id": str(tenant_id), "course_id": str(course_id)}
        ), collection_name

    async def _get_openai_embeddings(self, texts: List[str]) -> List[List[float]]:
        if not self.openai_client:
            raise RuntimeError("OpenAI client not configured - set OPENAI_API_KEY")
        all_embeddings = []
        for i in range(0, len(texts), 100):
            batch = texts[i:i+100]
            response = self.openai_client.embeddings.create(
                model=self.EMBEDDING_MODEL,
                input=batch
            )
            all_embeddings.extend([item.embedding for item in response.data])
        return all_embeddings

    def _extract_pdf_text(self, file_path: Path) -> List[Dict]:
        pages = []
        MAX_PAGES = 300

        # Method 1: PyMuPDF (fitz) — best quality extraction
        try:
            import fitz
            doc = fitz.open(str(file_path))
            total_pdf_pages = len(doc)
            logger.info(f"📄 PDF opened with PyMuPDF: {total_pdf_pages} pages")

            for page_num in range(min(total_pdf_pages, MAX_PAGES)):
                page = doc.load_page(page_num)

                # Try standard text extraction first
                text = page.get_text("text")

                # If standard fails, try blocks extraction (handles complex layouts)
                if not text.strip():
                    blocks = page.get_text("blocks")
                    text = "\n".join(
                        block[4] for block in blocks
                        if len(block) >= 5 and isinstance(block[4], str)
                    )

                if text.strip():
                    pages.append({"text": text.strip(), "page_number": page_num + 1})

            doc.close()

            if pages:
                logger.info(f"📄 PyMuPDF extracted text from {len(pages)}/{total_pdf_pages} pages")
                return pages

            # If fitz found no text on any page, this is likely a scanned/image PDF.
            # Try OCR via OpenAI Vision on a sample of pages.
            logger.warning(
                f"⚠️ PyMuPDF found 0 text pages in {total_pdf_pages}-page PDF. "
                f"Attempting OCR via OpenAI Vision on first 20 pages..."
            )
            if self.openai_client:
                import base64
                doc = fitz.open(str(file_path))
                ocr_pages = min(total_pdf_pages, 20)  # OCR first 20 pages max
                for page_num in range(ocr_pages):
                    try:
                        page = doc.load_page(page_num)
                        # Render page to image at 150 DPI (balance quality vs size)
                        pix = page.get_pixmap(dpi=150)
                        img_bytes = pix.tobytes("png")
                        img_b64 = base64.b64encode(img_bytes).decode("utf-8")

                        response = self.openai_client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=[
                                {"role": "system", "content": "Extract ALL text from this document page image. Return only the extracted text, no commentary."},
                                {"role": "user", "content": [
                                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}", "detail": "high"}}
                                ]}
                            ],
                            max_tokens=2000,
                            temperature=0.0,
                        )
                        ocr_text = response.choices[0].message.content.strip()
                        if ocr_text and len(ocr_text) > 20:
                            pages.append({"text": ocr_text, "page_number": page_num + 1})
                            logger.info(f"📄 OCR page {page_num + 1}: {len(ocr_text)} chars")
                    except Exception as ocr_e:
                        logger.warning(f"⚠️ OCR failed for page {page_num + 1}: {ocr_e}")
                doc.close()

                if pages:
                    logger.info(f"📄 OpenAI Vision OCR extracted text from {len(pages)} pages")
                    return pages

            logger.warning("⚠️ No text extracted even with OCR attempt")

        except ImportError:
            logger.warning("PyMuPDF not installed, falling back to PyPDF2")

        # Method 2: PyPDF2 fallback
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                total_pages = len(reader.pages)
                logger.info(f"📄 PyPDF2 reading {total_pages} pages")
                for i, page in enumerate(reader.pages[:MAX_PAGES]):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append({"text": text.strip(), "page_number": i + 1})
            if pages:
                logger.info(f"📄 PyPDF2 extracted text from {len(pages)} pages")
            return pages
        except Exception as e:
            logger.error(f"PDF extraction failed with both PyMuPDF and PyPDF2: {e}")
            raise ValueError(f"PDF extraction failed: {e}")

    def _extract_pptx_text(self, file_path: Path) -> List[Dict]:
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slides.append({"text": "\n".join(slide_text), "page_number": slide_num})
            return slides
        except ImportError:
            raise RuntimeError("Install python-pptx")
        except Exception as e:
            raise ValueError(f"PPTX extraction failed: {e}")

    def _extract_text_file(self, file_path: Path) -> List[Dict]:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        return [{"text": content, "page_number": 1}]

    def _extract_image_text(self, file_path: Path) -> List[Dict]:
        """Extract text from an image using OpenAI Vision (GPT-4o-mini)."""
        import base64

        if not self.openai_client:
            raise RuntimeError("OpenAI client required for image text extraction")

        ext = file_path.suffix.lower().lstrip('.')
        media_map = {'jpg': 'image/jpeg', 'jpeg': 'image/jpeg', 'png': 'image/png',
                     'webp': 'image/webp', 'bmp': 'image/bmp', 'gif': 'image/gif'}
        media_type = media_map.get(ext, 'image/png')

        with open(file_path, 'rb') as f:
            image_bytes = f.read()
        image_b64 = base64.b64encode(image_bytes).decode('utf-8')

        logger.info(f"🖼️ Extracting text from image via OpenAI Vision: {file_path.name} ({len(image_bytes)} bytes)")

        response = self.openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Extract ALL text, formulas, labels, and descriptions from this image. "
                 "If it is a diagram, chart, or illustration, describe its content in detail. "
                 "Return only the extracted/described content, no commentary."},
                {"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": f"data:{media_type};base64,{image_b64}", "detail": "high"}}
                ]}
            ],
            max_tokens=2000,
            temperature=0.0,
        )
        extracted_text = response.choices[0].message.content.strip()
        logger.info(f"🖼️ Image text extracted: {len(extracted_text)} chars")

        if not extracted_text or len(extracted_text) < 10:
            raise ValueError("No meaningful text could be extracted from the image")

        return [{"text": extracted_text, "page_number": 1}]

    def _extract_video_transcript(self, file_path: Path) -> List[Dict]:
        import tempfile
        audio_path = None
        try:
            from moviepy.editor import VideoFileClip
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
                audio_path = tmp.name
            clip = VideoFileClip(str(file_path))
            clip.audio.write_audiofile(audio_path, verbose=False, logger=None)
            clip.close()
            model = self._get_whisper_model()
            result = model.transcribe(audio_path, verbose=False)
            segments = []
            for seg in result.get("segments", []):
                segments.append({
                    "text": seg["text"].strip(),
                    "page_number": None,
                    "timestamp_start": seg["start"],
                    "timestamp_label": f"{int(seg['start'])//60}:{int(seg['start'])%60:02d}"
                })
            return segments
        finally:
            if audio_path and os.path.exists(audio_path):
                try:
                    os.unlink(audio_path)
                except Exception:
                    pass

    def _extract_youtube_transcript(self, youtube_url: str) -> List[Dict]:
        import tempfile
        audio_path = None
        try:
            import yt_dlp
            with tempfile.NamedTemporaryFile(suffix=".%(ext)s", delete=False) as tmp:
                tmp_template = tmp.name
            ydl_opts = {
                'format': 'bestaudio/best',
                'outtmpl': tmp_template,
                'postprocessors': [{'key': 'FFmpegExtractAudio', 'preferredcodec': 'wav'}],
                'quiet': True
            }
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.extract_info(youtube_url, download=True)
            audio_path = tmp_template.replace('%(ext)s', 'wav')
            if not os.path.exists(audio_path):
                for ext in ['m4a', 'webm', 'opus']:
                    p = tmp_template.replace('%(ext)s', ext)
                    if os.path.exists(p):
                        audio_path = p
                        break
            model = self._get_whisper_model()
            result = model.transcribe(audio_path, verbose=False)
            segments = []
            for seg in result.get("segments", []):
                segments.append({
                    "text": seg["text"].strip(),
                    "page_number": None,
                    "timestamp_start": seg["start"],
                    "timestamp_label": f"{int(seg['start'])//60}:{int(seg['start'])%60:02d}"
                })
            return segments
        except ImportError:
            raise RuntimeError("Install yt-dlp")
        finally:
            if audio_path and os.path.exists(audio_path):
                try:
                    os.unlink(audio_path)
                except Exception:
                    pass

    def _chunk_pages(self, pages: List[Dict], document_id: str, source_filename: str) -> List[Dict]:
        all_chunks = []
        for page in pages:
            text = page["text"]
            if not text.strip():
                continue
            if self.text_splitter:
                split_texts = self.text_splitter.split_text(text)
            else:
                split_texts = [text[i:i + self.CHUNK_SIZE * 4] for i in range(0, len(text), self.CHUNK_SIZE * 4)]
                split_texts = [t for t in split_texts if t.strip()]
            for split in split_texts:
                if not split.strip():
                    continue
                all_chunks.append({
                    "id": f"{document_id}_chunk_{len(all_chunks)}",
                    "text": split,
                    "source_file": source_filename,
                    "page_number": page.get("page_number"),
                    "timestamp_label": str(page.get("timestamp_label") or ""),
                    "timestamp_start": str(page.get("timestamp_start") or ""),
                })
        return all_chunks

    async def _index_to_chromadb(self, chunks: List[Dict], tenant_id: str, course_id: str,
                                 document_id: str, document_title: str, document_type: str) -> int:
        if not chunks or not self.chroma_client:
            return 0
        collection, _ = self._get_chroma_collection(tenant_id, course_id)
        texts = [c["text"] for c in chunks]
        embeddings = await self._get_openai_embeddings(texts)
        ids = [c["id"] for c in chunks]
        metadatas = []
        for c in chunks:
            metadatas.append({
                "document_id": document_id,
                "document_title": document_title,
                "document_type": document_type,
                "course_id": str(course_id),
                "tenant_id": str(tenant_id),
                "source_file": c.get("source_file", "unknown"),
                "page_number": str(c.get("page_number") or ""),
                "timestamp_label": str(c.get("timestamp_label") or ""),
                "timestamp_start": str(c.get("timestamp_start") or ""),
            })
        for i in range(0, len(ids), 500):
            collection.add(
                ids=ids[i:i+500],
                embeddings=embeddings[i:i+500],
                documents=texts[i:i+500],
                metadatas=metadatas[i:i+500]
            )
        return len(ids)

    async def index_document(
        self,
        file_content: bytes = None,
        filename: str = None,
        title: str = None,
        document_type: str = None,
        course_id: str = None,
        tenant_id: str = None,
        db_session=None,
        description: str = None,
        tags: List[str] = None,
        youtube_url: str = None
    ) -> Dict[str, Any]:
        from app.models.course_document import CourseDocument

        # Fail fast on excessively large uploads to avoid OOM/timeouts.
        if file_content is not None and len(file_content) > self.MAX_FILE_SIZE:
            raise ValueError(
                f"File too large: {len(file_content) / (1024 * 1024):.1f} MB. "
                f"Maximum allowed size is {self.MAX_FILE_SIZE / (1024 * 1024):.0f} MB."
            )

        document_id = str(uuid.uuid4())
        collection_name = f"tenant_{str(tenant_id)}_course_{str(course_id)}".replace("-", "_")
        doc = CourseDocument(
            id=document_id,
            title=title,
            document_type=document_type,
            status='processing',
            course_id=course_id,
            tenant_id=tenant_id,
            original_url=youtube_url,
            chroma_collection=collection_name,
            description=description or '',
            tags=tags or []
        )
        file_path = None
        if file_content and filename:
            safe_name = f"{document_id}_{filename}"
            fp = self.storage_path / str(tenant_id)
            fp.mkdir(parents=True, exist_ok=True)
            file_path = fp / safe_name
            with open(file_path, 'wb') as f:
                f.write(file_content)
            doc.file_path = str(file_path)
            doc.file_size = len(file_content)

        if db_session:
            db_session.add(doc)
            db_session.commit()

        try:
            pages = []
            source_filename = filename or youtube_url or title
            if document_type == 'pdf':
                pages = self._extract_pdf_text(file_path)
            elif document_type == 'pptx':
                pages = self._extract_pptx_text(file_path)
            elif document_type in ['text', 'markdown']:
                pages = self._extract_text_file(file_path)
            elif document_type == 'video':
                pages = self._extract_video_transcript(file_path)
            elif document_type == 'youtube':
                if not youtube_url:
                    raise ValueError("youtube_url required for YouTube type")
                pages = self._extract_youtube_transcript(youtube_url)
            elif document_type == 'image':
                pages = self._extract_image_text(file_path)
            else:
                raise ValueError(f"Unsupported document_type: {document_type}")

            if not pages:
                raise ValueError("No text extracted from document")

            # Protect against extremely long documents generating huge embeddings jobs.
            MAX_PAGES = 300
            if len(pages) > MAX_PAGES:
                logger.warning(
                    f"Truncating document {document_id}: {len(pages)} pages -> {MAX_PAGES} pages"
                )
                pages = pages[:MAX_PAGES]

            chunks = self._chunk_pages(pages, document_id, source_filename)
            vector_count = await self._index_to_chromadb(
                chunks, tenant_id, course_id, document_id, title, document_type
            )
            doc.status = 'indexed'
            doc.chunk_count = len(chunks)
            doc.vector_count = vector_count
            doc.last_updated = datetime.utcnow()
            doc.doc_metadata = {
                'page_count': max((p.get('page_number') or 0 for p in pages), default=0),
                'total_text_chars': sum(len(p['text']) for p in pages)
            }
            if db_session:
                db_session.commit()
            return {
                'document_id': document_id, 'title': title, 'document_type': document_type,
                'status': 'indexed', 'course_id': course_id, 'tenant_id': tenant_id,
                'upload_timestamp': datetime.utcnow(), 'file_size': len(file_content) if file_content else 0,
                'chunk_count': len(chunks), 'vector_count': vector_count,
                'processing_time': 0, 'chroma_collection': collection_name
            }
        except Exception as e:
            logger.error(f"Indexing failed: {e}")
            doc.status = 'error'
            doc.error_message = str(e)
            if db_session:
                db_session.commit()
            raise
